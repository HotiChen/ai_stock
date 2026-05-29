"""
make_pptx.py — 生成 QUANT·AI 系統流程 PowerPoint
執行：python3 make_pptx.py
輸出：QUANT_AI_流程圖.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.util as util

# ── 顏色定義 ─────────────────────────────────────────────────────────────────
BLACK      = RGBColor(0x0d, 0x11, 0x17)   # 背景
SURFACE    = RGBColor(0x16, 0x1b, 0x22)   # 卡片
UP         = RGBColor(0xC8, 0x33, 0x2B)   # 台股漲紅
DOWN       = RGBColor(0x2E, 0x7D, 0x4F)   # 台股跌綠
ACCENT     = RGBColor(0x26, 0x8B, 0xD2)   # 藍色重點
MUTED      = RGBColor(0x88, 0x88, 0x99)   # 灰色輔助
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
YELLOW     = RGBColor(0xD4, 0xA0, 0x17)
ORANGE     = RGBColor(0xF0, 0x8C, 0x00)
TEAL       = RGBColor(0x2A, 0x9D, 0x8F)
PURPLE     = RGBColor(0x8B, 0x5C, 0xF6)

# ── 投影片尺寸 16:9 ──────────────────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]   # completely blank


def new_slide():
    return prs.slides.add_slide(blank_layout)


def bg(slide, color=BLACK):
    """填滿背景色。"""
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1, 0, 0, W, H
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color


def box(slide, x, y, w, h, fill=SURFACE, line=None, radius=None):
    """加一個色塊。"""
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1.2)
    else:
        shape.line.fill.background()
    return shape


def txt(slide, text, x, y, w, h,
        size=18, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, wrap=True):
    """加文字方塊。"""
    tf = slide.shapes.add_textbox(x, y, w, h)
    tf.text_frame.word_wrap = wrap
    p  = tf.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.name  = "Helvetica Neue"
    return tf


def label_box(slide, title, subtitle, x, y, w, h, fill, accent):
    """小資訊卡片：彩色左邊框 + 標題 + 副標。"""
    box(slide, x, y, w, h, fill=SURFACE)
    # 左邊彩色條
    box(slide, x, y, Inches(0.08), h, fill=accent)
    txt(slide, title, x + Inches(0.18), y + Inches(0.08), w - Inches(0.25), Inches(0.35),
        size=14, bold=True, color=accent)
    txt(slide, subtitle, x + Inches(0.18), y + Inches(0.42), w - Inches(0.25), h - Inches(0.5),
        size=11, color=MUTED)


def arrow(slide, x1, y1, x2, y2, color=ACCENT):
    """水平箭頭（只支援水平線）。"""
    from pptx.util import Pt
    cx = (x1 + x2) // 2
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)
    connector.line.color.rgb = color
    connector.line.width = Pt(2)


# ════════════════════════════════════════════════════════════════════════════════
# Slide 1 — 封面
# ════════════════════════════════════════════════════════════════════════════════
s = new_slide()
bg(s)
# 橫線裝飾
box(s, 0, Inches(3.2), W, Inches(0.04), fill=UP)
box(s, 0, Inches(3.28), W, Inches(0.02), fill=ACCENT)

txt(s, "QUANT · AI", Inches(0.8), Inches(1.2), Inches(12), Inches(1.0),
    size=52, bold=True, color=UP, align=PP_ALIGN.CENTER)
txt(s, "台股 AI 量化交易系統", Inches(0.8), Inches(2.1), Inches(12), Inches(0.8),
    size=30, bold=False, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "完整運作流程 · 時間節點 · Telegram 推播", Inches(0.8), Inches(3.45), Inches(12), Inches(0.6),
    size=18, color=MUTED, align=PP_ALIGN.CENTER)

features = [
    ("🤖", "AI 選股", "Gemini + Claude"),
    ("📊", "當沖監控", "TP / SL / Trailing"),
    ("📱", "Telegram", "即時推播確認"),
    ("📈", "學習反饋", "越來越準"),
    ("🏦", "Notion 報告", "每日推送"),
]
fw = Inches(2.3)
fx = Inches(0.6)
for i, (icon, title, sub) in enumerate(features):
    bx = fx + i * (fw + Inches(0.1))
    box(s, bx, Inches(4.5), fw, Inches(1.6), fill=SURFACE)
    txt(s, icon,  bx, Inches(4.55), fw, Inches(0.5), size=22, align=PP_ALIGN.CENTER)
    txt(s, title, bx, Inches(5.05), fw, Inches(0.4), size=13, bold=True, color=WHITE,  align=PP_ALIGN.CENTER)
    txt(s, sub,   bx, Inches(5.45), fw, Inches(0.4), size=10, color=MUTED, align=PP_ALIGN.CENTER)

txt(s, "by QUANT·AI System  ·  2026", Inches(0), Inches(7.0), W, Inches(0.35),
    size=9, color=MUTED, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════════
# Slide 2 — 系統架構全覽
# ════════════════════════════════════════════════════════════════════════════════
s = new_slide()
bg(s)
txt(s, "系統架構全覽", Inches(0.5), Inches(0.25), Inches(10), Inches(0.6),
    size=28, bold=True, color=WHITE)
box(s, Inches(0.5), Inches(0.85), Inches(2.5), Inches(0.04), fill=UP)

layers = [
    ("既有 Python 後端\n（ai_stock/）", MUTED,   "選股 · 監控 · 下單 · 學習"),
    ("FastAPI 後端\n（port 8000）",    ACCENT,  "REST API + WebSocket 薄包裝層"),
    ("React 前端\n（port 5173）",      DOWN,    "Bloomberg 機構級 UI · 17 畫面"),
    ("Telegram Bot",                   YELLOW,  "即時確認 · 查股 · 緊急操作"),
    ("Notion 報告",                    PURPLE,  "每日學習報告 · 勝率趨勢"),
]
lw = Inches(2.3)
lx = Inches(0.35)
for i, (title, color, sub) in enumerate(layers):
    bx = lx + i * (lw + Inches(0.17))
    box(s, bx, Inches(1.2), lw, Inches(1.5), fill=SURFACE, line=color)
    txt(s, title, bx + Inches(0.15), Inches(1.3), lw - Inches(0.2), Inches(0.7),
        size=13, bold=True, color=color)
    txt(s, sub,   bx + Inches(0.15), Inches(2.0), lw - Inches(0.2), Inches(0.55),
        size=10, color=MUTED)

# 箭頭說明
txt(s, "↔  所有層均雙向通訊  ↔", Inches(0), Inches(2.85), W, Inches(0.35),
    size=11, color=MUTED, align=PP_ALIGN.CENTER)

# 下方模組列表
modules = [
    ("morning_briefing.py", "Gemini 早晨簡報"),
    ("stock_query.py", "技術指標 + DT 評分"),
    ("deep_analyzer.py", "Claude 深度分析"),
    ("risk_guard.py", "風控驗證"),
    ("monitor_agent.py", "Shioaji Tick 監控"),
    ("executor.py", "下單執行"),
    ("daytrading_review.py", "收盤複盤"),
    ("adaptive_scorer.py", "AI 學習更新"),
    ("notion_reporter.py", "Notion 推送"),
]
mw = Inches(1.37)
mx = Inches(0.35)
for i, (mod, desc) in enumerate(modules):
    row = i // 9
    col = i % 9
    bx = mx + col * (mw + Inches(0.05))
    by = Inches(3.3) + row * Inches(1.6)
    box(s, bx, by, mw, Inches(1.3), fill=SURFACE)
    txt(s, mod,  bx + Inches(0.08), by + Inches(0.08), mw - Inches(0.1), Inches(0.5),
        size=9, bold=True, color=ACCENT)
    txt(s, desc, bx + Inches(0.08), by + Inches(0.55), mw - Inches(0.1), Inches(0.6),
        size=9, color=WHITE)


# ════════════════════════════════════════════════════════════════════════════════
# Slide 3 — 每日完整時間軸
# ════════════════════════════════════════════════════════════════════════════════
s = new_slide()
bg(s)
txt(s, "每日完整時間軸", Inches(0.5), Inches(0.2), Inches(10), Inches(0.6),
    size=28, bold=True, color=WHITE)
box(s, Inches(0.5), Inches(0.78), Inches(2.5), Inches(0.04), fill=UP)

events = [
    ("08:30", "盤前 AI 分析",     ACCENT,   "選股 + 當沖預測\n→ Telegram 推播"),
    ("09:00", "開盤下單",          DOWN,     "波段下單執行"),
    ("09:05", "開盤確認",          TEAL,     "AI 再確認進場條件\n→ Telegram 開盤報告"),
    ("09:10", "當沖進場",          UP,       "買入確認 Telegram\n→ Shioaji 下單"),
    ("09:15", "Tick 監控開始",     YELLOW,   "MonitorAgent 訂閱 tick"),
    ("10:00", "盤中檢查",          ACCENT,   "持倉損益 + 候選提示\n→ Telegram 報告"),
    ("13:00", "出場警報",          ORANGE,   "損益預覽\n⚠️ 25 分鐘後強制平倉"),
    ("13:15", "停止 Tick 訂閱",   MUTED,    "MonitorAgent 停止"),
    ("13:25", "強制平倉",          UP,       "ForceCloseJob\n市價出清所有當沖部位"),
    ("13:35", "收盤任務",          PURPLE,   "PostMarketJob\nDT複盤 + 學習更新 + Notion"),
]

# 時間軸橫線
ty = Inches(4.3)
box(s, Inches(0.5), ty, Inches(12.33), Inches(0.04), fill=ACCENT)

ew = Inches(1.18)
ex = Inches(0.45)
for i, (time_str, title, color, detail) in enumerate(events):
    bx = ex + i * (ew + Inches(0.05))
    # 圓點
    dot = s.shapes.add_shape(9, bx + ew//2 - Inches(0.12), ty - Inches(0.14), Inches(0.28), Inches(0.28))
    dot.fill.solid(); dot.fill.fore_color.rgb = color
    dot.line.fill.background()
    # 時間標籤（交錯上下）
    if i % 2 == 0:
        txt(s, time_str, bx, ty - Inches(0.55), ew, Inches(0.4),
            size=13, bold=True, color=color, align=PP_ALIGN.CENTER)
        box(s, bx + Inches(0.05), ty + Inches(0.35), ew - Inches(0.1), Inches(1.85), fill=SURFACE)
        txt(s, title,  bx + Inches(0.1), ty + Inches(0.4),  ew - Inches(0.15), Inches(0.45),
            size=11, bold=True, color=color)
        txt(s, detail, bx + Inches(0.1), ty + Inches(0.85), ew - Inches(0.15), Inches(1.2),
            size=9, color=MUTED)
    else:
        txt(s, time_str, bx, ty + Inches(0.1), ew, Inches(0.4),
            size=13, bold=True, color=color, align=PP_ALIGN.CENTER)
        box(s, bx + Inches(0.05), ty - Inches(2.3), ew - Inches(0.1), Inches(1.85), fill=SURFACE)
        txt(s, title,  bx + Inches(0.1), ty - Inches(2.25), ew - Inches(0.15), Inches(0.45),
            size=11, bold=True, color=color)
        txt(s, detail, bx + Inches(0.1), ty - Inches(1.8),  ew - Inches(0.15), Inches(1.2),
            size=9, color=MUTED)


# ════════════════════════════════════════════════════════════════════════════════
# Slide 4 — 08:30 盤前 AI 分析流程
# ════════════════════════════════════════════════════════════════════════════════
s = new_slide()
bg(s)
box(s, 0, 0, W, Inches(1.0), fill=SURFACE)
txt(s, "08:30", Inches(0.4), Inches(0.15), Inches(1.5), Inches(0.6),
    size=32, bold=True, color=UP)
txt(s, "盤前 AI 分析 — PremarketJob + DaytradingReport", Inches(2.0), Inches(0.22), Inches(11), Inches(0.55),
    size=20, bold=True, color=WHITE)

steps = [
    ("① Gemini 早晨簡報", ACCENT,
     "• 抓取台股加權 / 外資 / 投信 / SOX / VIX\n• Gemini 生成市場分析 + 風險提示\n• 輸出 MorningBriefing 物件"),
    ("② risk_guard 風控", YELLOW,
     "• 檢查每筆預算 ≤ 30%\n• 同板塊持倉上限\n• 黑名單過濾"),
    ("③ 三套策略生成", DOWN,
     "• 衝刺版（高風高報）\n• 均衡版（平衡）\n• 保守版（防守）\n• 每套含選股 + 預期報酬"),
    ("④ 當沖預測報告", PURPLE,
     "• 掃描全市場\n• dt_score 技術評分（0–10）\n• Haiku AI 分析評分≥4的標的\n• 存入 daytrading_review.db"),
    ("⑤ Telegram 推播", ORANGE,
     "• 早晨市場簡報\n• 三套策略各一則\n• Inline Keyboard：\n  [🚀衝刺] [⚖️均衡] [🛡️保守]"),
]
sw = Inches(2.42)
sx = Inches(0.3)
for i, (title, color, detail) in enumerate(steps):
    bx = sx + i * (sw + Inches(0.08))
    box(s, bx, Inches(1.2), sw, Inches(5.7), fill=SURFACE, line=color)
    # 頂部色條
    box(s, bx, Inches(1.2), sw, Inches(0.06), fill=color)
    txt(s, title, bx + Inches(0.12), Inches(1.35), sw - Inches(0.2), Inches(0.5),
        size=14, bold=True, color=color)
    txt(s, detail, bx + Inches(0.12), Inches(1.9), sw - Inches(0.2), Inches(4.8),
        size=11, color=WHITE)

# 底部備注
txt(s, "使用者收到 Telegram 後點選策略 → 系統記錄選擇 → 09:00 執行", Inches(0), Inches(7.05), W, Inches(0.35),
    size=10, color=MUTED, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════════
# Slide 5 — 09:00–09:10 開盤進場流程
# ════════════════════════════════════════════════════════════════════════════════
s = new_slide()
bg(s)
box(s, 0, 0, W, Inches(1.0), fill=SURFACE)
txt(s, "09:00 – 09:10", Inches(0.4), Inches(0.15), Inches(3.5), Inches(0.6),
    size=28, bold=True, color=DOWN)
txt(s, "開盤進場三步驟", Inches(4.0), Inches(0.22), Inches(9), Inches(0.55),
    size=20, bold=True, color=WHITE)

phase_x = [Inches(0.4), Inches(4.55), Inches(8.75)]
phase_w = Inches(3.9)

# 09:00
box(s, phase_x[0], Inches(1.1), phase_w, Inches(5.9), fill=SURFACE)
box(s, phase_x[0], Inches(1.1), phase_w, Inches(0.07), fill=DOWN)
txt(s, "09:00  波段下單", phase_x[0] + Inches(0.15), Inches(1.25), phase_w, Inches(0.5),
    size=16, bold=True, color=DOWN)
txt(s, "MarketOpenJob",   phase_x[0] + Inches(0.15), Inches(1.75), phase_w, Inches(0.35),
    size=10, color=MUTED)
items_900 = [
    "✅ 使用者已選好策略",
    "→ 執行 picks（open/add/close）",
    "→ place_stock_order() via Shioaji",
    "→ 寫入 research.db",
    "",
    "⚠️ 每筆送出前：",
    "  send_confirmation() 二次確認",
    "  60 秒 timeout → 自動略過",
]
for j, item in enumerate(items_900):
    txt(s, item, phase_x[0] + Inches(0.15), Inches(2.15) + j * Inches(0.5),
        phase_w - Inches(0.2), Inches(0.45), size=11, color=WHITE if item else MUTED)

# 09:05
box(s, phase_x[1], Inches(1.1), phase_w, Inches(5.9), fill=SURFACE)
box(s, phase_x[1], Inches(1.1), phase_w, Inches(0.07), fill=TEAL)
txt(s, "09:05  開盤確認", phase_x[1] + Inches(0.15), Inches(1.25), phase_w, Inches(0.5),
    size=16, bold=True, color=TEAL)
txt(s, "run_opening_reconfirm()",phase_x[1] + Inches(0.15), Inches(1.75), phase_w, Inches(0.35),
    size=10, color=MUTED)
items_905 = [
    "AI 再次確認進場條件：",
    "  • 量能比 ≥ 0.8？",
    "  • 大盤方向吻合？",
    "  • 台指期溢貼水正常？",
    "",
    "→ confirmed：繼續",
    "→ skipped：放棄，通知使用者",
    "",
    "📱 Telegram 推播：",
    "  ✅/❌ 每支股狀態",
    "  現價 + 進場區間",
]
for j, item in enumerate(items_905):
    txt(s, item, phase_x[1] + Inches(0.15), Inches(2.15) + j * Inches(0.42),
        phase_w - Inches(0.2), Inches(0.4), size=11, color=WHITE)

# 09:10
box(s, phase_x[2], Inches(1.1), phase_w, Inches(5.9), fill=SURFACE)
box(s, phase_x[2], Inches(1.1), phase_w, Inches(0.07), fill=UP)
txt(s, "09:10  當沖進場", phase_x[2] + Inches(0.15), Inches(1.25), phase_w, Inches(0.5),
    size=16, bold=True, color=UP)
txt(s, "send_dt_buy_confirmation()", phase_x[2] + Inches(0.15), Inches(1.75), phase_w, Inches(0.35),
    size=10, color=MUTED)
items_910 = [
    "📱 Telegram 發出：",
    "",
    "  [✅ 2330 台積電] [❌ 跳過]",
    "  [✅ 2382 廣達]   [❌ 跳過]",
    "  [✅ 全部買入]    [❌ 全部跳過]",
    "",
    "使用者點選 → callback 觸發",
    "→ _handle_dt_buy(code)",
    "→ Shioaji 下買單",
    "→ MonitorAgent 開始 tick 監控",
    "→ 持倉狀態 watching → active",
]
for j, item in enumerate(items_910):
    color = ACCENT if item.startswith("  [") else WHITE
    txt(s, item, phase_x[2] + Inches(0.15), Inches(2.15) + j * Inches(0.42),
        phase_w - Inches(0.2), Inches(0.4), size=11, color=color)


# ════════════════════════════════════════════════════════════════════════════════
# Slide 6 — 盤中 Tick 監控 + 即時賣出邏輯
# ════════════════════════════════════════════════════════════════════════════════
s = new_slide()
bg(s)
box(s, 0, 0, W, Inches(1.0), fill=SURFACE)
txt(s, "09:15 – 13:15", Inches(0.4), Inches(0.15), Inches(3.5), Inches(0.6),
    size=28, bold=True, color=YELLOW)
txt(s, "Tick 監控 + 即時出場決策", Inches(4.0), Inches(0.22), Inches(9), Inches(0.55),
    size=20, bold=True, color=WHITE)

# 左側：監控架構
box(s, Inches(0.3), Inches(1.1), Inches(4.0), Inches(5.9), fill=SURFACE)
txt(s, "MonitorAgent", Inches(0.45), Inches(1.2), Inches(3.8), Inches(0.5),
    size=16, bold=True, color=YELLOW)
arch_items = [
    ("訂閱 Shioaji Tick", ACCENT),
    ("每個 tick 呼叫 check_trailing_stop()", WHITE),
    ("", WHITE),
    ("check_trailing_stop() 判斷優先順序：", MUTED),
    ("", WHITE),
    ("① gain_pct ≤ -3.0%  →  停損出場", UP),
    ("② gain_pct ≥ +9.0%  →  達目標出場", DOWN),
    ("③ 時間 ≥ 13:00  →  強制出場", ORANGE),
    ("④ peak_gain ≥ +3% 且", YELLOW),
    ("   從高點回落 ≥ 2%  →  追蹤停損", YELLOW),
    ("", WHITE),
    ("peak_price 每個 tick 持續更新", MUTED),
    ("（只漲不跌）", MUTED),
]
for j, (item, color) in enumerate(arch_items):
    txt(s, item, Inches(0.5), Inches(1.8) + j * Inches(0.38),
        Inches(3.7), Inches(0.35), size=10, color=color)

# 右側：觸發流程
box(s, Inches(4.6), Inches(1.1), Inches(8.4), Inches(5.9), fill=BLACK)

triggers = [
    ("停損出場", UP, "3.0%", "gain_pct ≤ -stop_loss_pct",
     "立即市價賣出\n損失最小化"),
    ("達目標出場", DOWN, "9.0%", "gain_pct ≥ take_profit_pct",
     "鎖住最大獲利\n送出市價單"),
    ("追蹤停損", YELLOW, "3%/2%", "peak_gain ≥ trailing_start\ngap ≥ trailing_gap",
     "從高點保護利潤\n不急著出場"),
    ("強制平倉", ORANGE, "13:00", "當前時間 ≥ force_close_time",
     "確保當沖不留倉\n全數平倉"),
]
tw = Inches(1.9)
tx = Inches(4.75)
for i, (title, color, threshold, condition, action) in enumerate(triggers):
    bx = tx + i * (tw + Inches(0.12))
    box(s, bx, Inches(1.2), tw, Inches(5.6), fill=SURFACE, line=color)
    box(s, bx, Inches(1.2), tw, Inches(0.06), fill=color)
    txt(s, title,     bx + Inches(0.1), Inches(1.35), tw, Inches(0.4),
        size=13, bold=True, color=color)
    txt(s, "閾值：" + threshold, bx + Inches(0.1), Inches(1.8), tw, Inches(0.35),
        size=10, color=ACCENT)
    txt(s, "條件：", bx + Inches(0.1), Inches(2.25), tw, Inches(0.3), size=9, color=MUTED)
    txt(s, condition, bx + Inches(0.1), Inches(2.55), tw, Inches(0.7),
        size=9, color=WHITE)
    txt(s, "動作：", bx + Inches(0.1), Inches(3.35), tw, Inches(0.3), size=9, color=MUTED)
    txt(s, action,    bx + Inches(0.1), Inches(3.65), tw, Inches(0.6),
        size=9, color=DOWN)
    txt(s, "📱 Telegram 即時通知",  bx + Inches(0.1), Inches(4.4),  tw, Inches(0.35),
        size=9, color=YELLOW)
    txt(s, "✅ 出場已送出\n  + 觸發原因 + 價格", bx + Inches(0.1), Inches(4.75), tw, Inches(0.7),
        size=9, color=MUTED)


# ════════════════════════════════════════════════════════════════════════════════
# Slide 7 — 13:00 / 13:25 / 13:35 收盤三連發
# ════════════════════════════════════════════════════════════════════════════════
s = new_slide()
bg(s)
box(s, 0, 0, W, Inches(1.0), fill=SURFACE)
txt(s, "13:00 – 13:35", Inches(0.4), Inches(0.15), Inches(3.5), Inches(0.6),
    size=28, bold=True, color=ORANGE)
txt(s, "收盤三連發 — 警報 · 平倉 · 複盤", Inches(4.0), Inches(0.22), Inches(9), Inches(0.55),
    size=20, bold=True, color=WHITE)

phase2_x = [Inches(0.3), Inches(4.55), Inches(8.75)]

# 13:00
box(s, phase2_x[0], Inches(1.1), Inches(4.0), Inches(5.9), fill=SURFACE, line=ORANGE)
box(s, phase2_x[0], Inches(1.1), Inches(4.0), Inches(0.07), fill=ORANGE)
txt(s, "13:00  出場警報", phase2_x[0] + Inches(0.15), Inches(1.25), Inches(3.8), Inches(0.5),
    size=16, bold=True, color=ORANGE)
items_1300 = [
    "📱 Telegram 推播內容：",
    "",
    "📊 當沖出場警報 13:00",
    "🔴 2330 台積電",
    "  買 1068 → 現 1082",
    "  損益 +14,000（+1.3%）",
    "🔴 2382 廣達",
    "  買 305 → 現 312",
    "  損益 +14,000（+2.3%）",
    "",
    "💰 合計損益：+28,000 元",
    "⚠️ 13:25 將自動強制平倉",
]
for j, item in enumerate(items_1300):
    color = ORANGE if "警報" in item else (MUTED if item.startswith("  ") else WHITE)
    if "⚠️" in item:
        color = UP
    txt(s, item, phase2_x[0] + Inches(0.15), Inches(1.85) + j * Inches(0.38),
        Inches(3.7), Inches(0.35), size=10, color=color)

# 13:25
box(s, phase2_x[1], Inches(1.1), Inches(4.0), Inches(5.9), fill=SURFACE, line=UP)
box(s, phase2_x[1], Inches(1.1), Inches(4.0), Inches(0.07), fill=UP)
txt(s, "13:25  強制平倉", phase2_x[1] + Inches(0.15), Inches(1.25), Inches(3.8), Inches(0.5),
    size=16, bold=True, color=UP)
txt(s, "ForceCloseJob", phase2_x[1] + Inches(0.15), Inches(1.75), Inches(3.8), Inches(0.35),
    size=10, color=MUTED)
items_1325 = [
    "掃描所有 active 持倉",
    "",
    "Real 模式：",
    "→ place_order(price=0, MKT)",
    "→ 市價單確保成交",
    "→ 等待 fill 回報",
    "→ 寫入 force_close_requested",
    "",
    "Sim 模式：",
    "→ 取即時報價計算損益",
    "→ 直接寫入已平倉記錄",
    "",
    "100% 確保不留倉過夜",
]
for j, item in enumerate(items_1325):
    color = DOWN if item.startswith("→") else (ACCENT if "Sim" in item or "Real" in item else WHITE)
    txt(s, item, phase2_x[1] + Inches(0.15), Inches(2.15) + j * Inches(0.38),
        Inches(3.7), Inches(0.35), size=10, color=color)

# 13:35
box(s, phase2_x[2], Inches(1.1), Inches(4.0), Inches(5.9), fill=SURFACE, line=PURPLE)
box(s, phase2_x[2], Inches(1.1), Inches(4.0), Inches(0.07), fill=PURPLE)
txt(s, "13:35  收盤任務", phase2_x[2] + Inches(0.15), Inches(1.25), Inches(3.8), Inches(0.5),
    size=16, bold=True, color=PURPLE)
items_1335 = [
    "① PostMarketJob",
    "  停止監控 + 存每日摘要",
    "  notify_market_close()",
    "",
    "② DaytradingReview",
    "  回填 OHLC（yfinance）",
    "  判斷 hit_target / hit_stop",
    "  → Telegram 複盤報告",
    "",
    "③ adaptive_scorer.py",
    "  更新評分區間勝率",
    "  生成 learning_summary",
    "  寫入 adaptive_config.json",
    "",
    "④ notion_reporter.py",
    "  推送今日完整報告",
    "  → Notion 資料庫",
]
for j, item in enumerate(items_1335):
    color = MUTED if item.startswith("  ") else WHITE
    if item[0:1] in ["①","②","③","④"]:
        color = YELLOW
    txt(s, item, phase2_x[2] + Inches(0.15), Inches(1.85) + j * Inches(0.33),
        Inches(3.7), Inches(0.32), size=10, color=color)


# ════════════════════════════════════════════════════════════════════════════════
# Slide 8 — Telegram 推播完整清單
# ════════════════════════════════════════════════════════════════════════════════
s = new_slide()
bg(s)
txt(s, "Telegram 推播完整清單", Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
    size=28, bold=True, color=WHITE)
box(s, Inches(0.5), Inches(0.78), Inches(3.5), Inches(0.04), fill=YELLOW)

messages = [
    ("系統啟動",  MUTED,   "notify_system_start()",           "系統已啟動 · 模擬/真實模式"),
    ("08:30",    ACCENT,  "send_morning_push()",             "📊 早晨市場簡報 + 三套策略 + 選擇 Keyboard"),
    ("08:30",    PURPLE,  "DaytradingReport",                "🎯 今日當沖候選股 · 評分 · 信心度"),
    ("09:05",    TEAL,    "開盤確認通知",                      "⚡ 每支股 ✅確認/❌放棄 + 現價 + 大盤方向"),
    ("09:10",    UP,      "send_dt_buy_confirmation()",      "💰 [✅ CODE] [❌跳過] Inline Keyboard"),
    ("盤中即時",  YELLOW,  "_run_dt_sell_alerts()",           "✅出場已送出 / ❌出場失敗 + 觸發原因 + 價格"),
    ("10:00",    ACCENT,  "_mid_session_check()",            "📊 持倉損益 + 剩餘候選提示"),
    ("13:00",    ORANGE,  "出場警報",                          "所有持倉損益 + ⚠️ 13:25 強制平倉"),
    ("13:25",    UP,      "ForceCloseJob",                   "強制平倉執行結果通知"),
    ("13:35",    DOWN,    "notify_market_close()",           "📈 今日總損益 + 交易筆數"),
    ("13:35",    PURPLE,  "DaytradingReview → send_text()", "🎯 預測複盤：hit_target / hit_stop · 近30日勝率"),
    ("系統停止",  MUTED,   "notify_system_stop()",           "系統已停止"),
]

# 表頭
hx = [Inches(0.3), Inches(1.5), Inches(4.0), Inches(7.2)]
for header, hbx in zip(["時間", "顏色分類", "函式 / 模組", "推播內容"], hx):
    txt(s, header, hbx, Inches(0.9), Inches(2.5), Inches(0.4),
        size=10, bold=True, color=MUTED)
box(s, Inches(0.3), Inches(1.3), Inches(12.7), Inches(0.03), fill=MUTED)

for i, (time_str, color, func, content) in enumerate(messages):
    row_y = Inches(1.4) + i * Inches(0.47)
    if i % 2 == 0:
        box(s, Inches(0.3), row_y, Inches(12.7), Inches(0.45), fill=SURFACE)
    txt(s, time_str, hx[0], row_y + Inches(0.05), Inches(1.1), Inches(0.38),
        size=11, bold=True, color=color)
    dot2 = s.shapes.add_shape(9, hx[1] + Inches(0.6), row_y + Inches(0.14), Inches(0.18), Inches(0.18))
    dot2.fill.solid(); dot2.fill.fore_color.rgb = color; dot2.line.fill.background()
    txt(s, func,    hx[1] + Inches(0.85), row_y + Inches(0.05), Inches(3.1), Inches(0.38),
        size=10, color=WHITE)
    txt(s, content, hx[2] + Inches(0.1),  row_y + Inches(0.05), Inches(5.9), Inches(0.38),
        size=10, color=MUTED)


# ════════════════════════════════════════════════════════════════════════════════
# Slide 9 — AI 學習反饋回路
# ════════════════════════════════════════════════════════════════════════════════
s = new_slide()
bg(s)
txt(s, "AI 學習反饋回路 — 越來越準", Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
    size=28, bold=True, color=WHITE)
box(s, Inches(0.5), Inches(0.78), Inches(3.0), Inches(0.04), fill=PURPLE)

# 環形流程圖（用方塊 + 箭頭模擬）
cycle_items = [
    ("① 盤前 AI 分析\n08:30", ACCENT,  Inches(0.4),  Inches(1.3)),
    ("② 使用者執行\n09:10 買入", DOWN,  Inches(4.8),  Inches(1.3)),
    ("③ 盤中監控\nTP/SL 觸發", YELLOW, Inches(9.2),  Inches(1.3)),
    ("④ 收盤複盤\n13:35 Review", ORANGE,Inches(9.2),  Inches(4.0)),
    ("⑤ 學習更新\nadaptive_scorer", PURPLE,Inches(4.8),Inches(4.0)),
    ("⑥ Notion 報告\n每日推送", TEAL,  Inches(0.4),  Inches(4.0)),
]
cw = Inches(3.4)
ch = Inches(1.8)
for title, color, cx, cy in cycle_items:
    box(s, cx, cy, cw, ch, fill=SURFACE, line=color)
    box(s, cx, cy, cw, Inches(0.07), fill=color)
    txt(s, title, cx + Inches(0.15), cy + Inches(0.2), cw - Inches(0.2), ch - Inches(0.3),
        size=14, bold=True, color=color)

# 連接箭頭說明
arrow_labels = [
    (Inches(3.8),  Inches(2.2), "評分預測"),
    (Inches(8.4),  Inches(2.2), "成交記錄"),
    (Inches(9.9),  Inches(3.7), "OHLC 回填"),
    (Inches(8.4),  Inches(4.9), "勝率統計"),
    (Inches(3.8),  Inches(4.9), "learning_summary"),
    (Inches(0.95), Inches(3.7), "注入隔日提示"),
]
for lx, ly, label in arrow_labels:
    txt(s, "→ " + label, lx, ly, Inches(2.0), Inches(0.35), size=9, color=MUTED)

# 右下：學習效果說明
box(s, Inches(0.3), Inches(6.1), Inches(12.5), Inches(1.2), fill=SURFACE)
txt(s, "學習效果（隨交易日累積）：",
    Inches(0.5), Inches(6.2), Inches(4), Inches(0.4), size=12, bold=True, color=YELLOW)
effect_items = [
    "評分 4-5 分  →  歷史勝率 31%  →  建議跳過",
    "評分 6-7 分  →  歷史勝率 58%  →  推薦操作",
    "評分 8-10分  →  歷史勝率 72%  →  強烈推薦",
]
for j, item in enumerate(effect_items):
    txt(s, item, Inches(0.5) + j * Inches(4.2), Inches(6.6), Inches(4.0), Inches(0.5),
        size=10, color=WHITE)
txt(s, "→ optimal_min_score 自動調高門檻，每日 Notion 顯示勝率趨勢",
    Inches(0.5), Inches(7.1), Inches(12), Inches(0.35), size=10, color=PURPLE)


# ════════════════════════════════════════════════════════════════════════════════
# Slide 10 — Notion 每日報告內容
# ════════════════════════════════════════════════════════════════════════════════
s = new_slide()
bg(s)
txt(s, "Notion 每日報告 — 📈 當沖 AI 每日報告", Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
    size=24, bold=True, color=WHITE)
box(s, Inches(0.5), Inches(0.78), Inches(3.0), Inches(0.04), fill=TEAL)

# 左：資料庫欄位
box(s, Inches(0.3), Inches(1.1), Inches(5.5), Inches(5.9), fill=SURFACE)
txt(s, "資料庫屬性欄位", Inches(0.45), Inches(1.2), Inches(5.2), Inches(0.45),
    size=14, bold=True, color=TEAL)
fields = [
    ("日期",      "Title",  "2026-05-29"),
    ("勝率%",     "Number", "58.3%"),
    ("總損益",    "Number", "+28,400 NTD"),
    ("預測筆數",  "Number", "6"),
    ("達目標",    "Number", "3"),
    ("觸停損",    "Number", "2"),
    ("未觸發",    "Number", "1"),
    ("大盤漲跌%", "Number", "+1.24%"),
    ("近7日勝率%","Number", "55.1%"),
    ("近30日勝率%","Number","52.8%"),
    ("最高分股票","Text",   "2330 台積電（9/10）"),
    ("結果",      "Select", "獲利日 / 虧損日 / 平盤"),
    ("學習版本",  "Number", "v12"),
    ("AI建議調整","Text",   "learning_summary（2000字）"),
]
for j, (name, ftype, example) in enumerate(fields):
    fy = Inches(1.75) + j * Inches(0.36)
    txt(s, name,    Inches(0.5), fy, Inches(1.8), Inches(0.32), size=10, bold=True, color=WHITE)
    txt(s, ftype,   Inches(2.3), fy, Inches(1.0), Inches(0.32), size=9, color=ACCENT)
    txt(s, example, Inches(3.3), fy, Inches(2.3), Inches(0.32), size=9, color=MUTED)

# 右：頁面內容
box(s, Inches(6.1), Inches(1.1), Inches(6.9), Inches(5.9), fill=SURFACE)
txt(s, "頁面正文（Markdown 格式）", Inches(6.25), Inches(1.2), Inches(6.5), Inches(0.45),
    size=14, bold=True, color=TEAL)
content_items = [
    ("## 2026-05-29 當沖複盤", WHITE),
    ("大盤漲跌 +1.24%　今日損益 +28,400 元", MUTED),
    ("", WHITE),
    ("## 預測結果（共 6 筆）", WHITE),
    ("✅ 達目標：3 筆", DOWN),
    ("❌ 觸停損：2 筆", UP),
    ("⬜ 未觸發：1 筆", MUTED),
    ("今日勝率：50.0%", WHITE),
    ("", WHITE),
    ("## 個股明細", WHITE),
    ("✅ 2330 台積電（評分 9/10）— 達目標", DOWN),
    ("   目標 1120  停損 1030  高 1125...", MUTED),
    ("❌ 2382 廣達（評分 7/10）— 觸停損", UP),
    ("", WHITE),
    ("## 🧠 AI 學習摘要", PURPLE),
    ("評分 8-10 歷史勝率：72.0%...", MUTED),
    ("## 各評分區間勝率（累積）", WHITE),
    ("評分 8-10：████████████░░░░░░░░ 72.0%", YELLOW),
    ("評分 6-7： ███████████░░░░░░░░░ 58.0%", YELLOW),
]
for j, (item, color) in enumerate(content_items):
    txt(s, item, Inches(6.25), Inches(1.75) + j * Inches(0.32),
        Inches(6.5), Inches(0.3), size=9, color=color)


# ════════════════════════════════════════════════════════════════════════════════
# Slide 11 — 環境設定 & 快速啟動
# ════════════════════════════════════════════════════════════════════════════════
s = new_slide()
bg(s)
txt(s, "環境設定 & 快速啟動", Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
    size=28, bold=True, color=WHITE)
box(s, Inches(0.5), Inches(0.78), Inches(2.0), Inches(0.04), fill=DOWN)

# .env 設定
box(s, Inches(0.3), Inches(1.1), Inches(6.0), Inches(5.9), fill=SURFACE)
txt(s, ".env 必填設定", Inches(0.45), Inches(1.2), Inches(5.5), Inches(0.45),
    size=14, bold=True, color=DOWN)
env_items = [
    ("# Telegram", MUTED),
    ("TELEGRAM_BOT_TOKEN=你的 Bot Token", ACCENT),
    ("TELEGRAM_CHAT_ID=你的 Chat ID", ACCENT),
    ("TELEGRAM_USER_ID=你的 User ID", ACCENT),
    ("", WHITE),
    ("# Notion（每日推報告）", MUTED),
    ("NOTION_TOKEN=secret_xxx", PURPLE),
    ("NOTION_DB_ID=1031a6aff6b3...  # 可選", MUTED),
    ("", WHITE),
    ("# 券商 API", MUTED),
    ("SHIOAJI_API_KEY=xxx", YELLOW),
    ("SHIOAJI_SECRET_KEY=xxx", YELLOW),
    ("SHIOAJI_SIMULATION=true", YELLOW),
    ("", WHITE),
    ("# AI 模型", MUTED),
    ("ANTHROPIC_API_KEY=sk-ant-xxx", TEAL),
    ("GEMINI_API_KEY=xxx", TEAL),
    ("", WHITE),
    ("# 下單安全上限", MUTED),
    ("ORDER_HARD_LIMIT=150000", WHITE),
    ("BUDGET_PER_STOCK=100000", WHITE),
]
for j, (item, color) in enumerate(env_items):
    txt(s, item, Inches(0.5), Inches(1.75) + j * Inches(0.28),
        Inches(5.5), Inches(0.26), size=9, color=color)

# 右：啟動指令
box(s, Inches(6.6), Inches(1.1), Inches(6.4), Inches(2.8), fill=SURFACE)
txt(s, "系統啟動", Inches(6.75), Inches(1.2), Inches(6.0), Inches(0.45),
    size=14, bold=True, color=DOWN)
start_items = [
    ("# 主排程（每日 08:30–13:35）", MUTED),
    ("python3 main.py", DOWN),
    ("", WHITE),
    ("# Telegram Bot（另開終端）", MUTED),
    ("python3 telegram_bot.py", YELLOW),
    ("", WHITE),
    ("# 每日 13:35 後手動執行", MUTED),
    ("bash post_market.sh", PURPLE),
    ("# 或加入 crontab：", MUTED),
    ("35 13 * * 1-5  bash post_market.sh", PURPLE),
]
for j, (item, color) in enumerate(start_items):
    txt(s, item, Inches(6.75), Inches(1.75) + j * Inches(0.3),
        Inches(6.0), Inches(0.28), size=9, color=color)

box(s, Inches(6.6), Inches(4.1), Inches(6.4), Inches(2.8), fill=SURFACE)
txt(s, "手動操作指令", Inches(6.75), Inches(4.2), Inches(6.0), Inches(0.45),
    size=14, bold=True, color=ACCENT)
manual_items = [
    ("# 查看當前學習狀態", MUTED),
    ("python3 adaptive_scorer.py", ACCENT),
    ("", WHITE),
    ("# 手動推送今日 Notion 報告", MUTED),
    ("python3 notion_reporter.py 2026-05-29", TEAL),
    ("", WHITE),
    ("# 當沖預測報告", MUTED),
    ("python3 daytrading_report.py", YELLOW),
    ("", WHITE),
    ("# Telegram Bot 互動（主選單）", MUTED),
    ("發送：/menu 到 Bot", WHITE),
]
for j, (item, color) in enumerate(manual_items):
    txt(s, item, Inches(6.75), Inches(4.65) + j * Inches(0.27),
        Inches(6.0), Inches(0.25), size=9, color=color)


# ════════════════════════════════════════════════════════════════════════════════
# Slide 12 — 總結
# ════════════════════════════════════════════════════════════════════════════════
s = new_slide()
bg(s)
box(s, 0, Inches(3.0), W, Inches(0.04), fill=UP)
box(s, 0, Inches(3.07), W, Inches(0.02), fill=ACCENT)

txt(s, "系統設計亮點", Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
    size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

highlights = [
    ("買入手動確認\n賣出全自動", UP,
     "下單需 Telegram 二次確認\n出場不等人，即時執行"),
    ("四層出場保護", YELLOW,
     "停損 3% / 達目標 9%\n追蹤停損 / 13:25 強制平倉"),
    ("AI 每日學習", PURPLE,
     "評分區間勝率統計\n自動調整最低門檻"),
    ("Notion 報告", TEAL,
     "每日自動推送\n勝率趨勢可視化"),
    ("Telegram 全程追蹤", ACCENT,
     "12 個推播節點\n系統啟動到收盤全覆蓋"),
    ("台股色彩規範", DOWN,
     "漲紅 #C8332B\n跌綠 #2E7D4F"),
]
hw = Inches(2.0)
hx2 = Inches(0.4)
for i, (title, color, desc) in enumerate(highlights):
    col = i % 3
    row = i // 3
    bx = hx2 + col * (hw + Inches(0.2))
    by = Inches(3.4) + row * Inches(1.8)
    box(s, bx, by, hw + Inches(0.1), Inches(1.65), fill=SURFACE, line=color)
    txt(s, title, bx + Inches(0.12), by + Inches(0.15), hw, Inches(0.65),
        size=13, bold=True, color=color)
    txt(s, desc,  bx + Inches(0.12), by + Inches(0.82), hw, Inches(0.7),
        size=10, color=MUTED)

# 右下大字
txt(s, "所有 17 個畫面 ·\n12 個 Telegram 推播節點 ·\n全自動排程", Inches(6.5), Inches(3.5), Inches(6.6), Inches(3.5),
    size=22, bold=True, color=WHITE)
box(s, Inches(6.5), Inches(6.8), Inches(6.5), Inches(0.04), fill=MUTED)
txt(s, "QUANT · AI · 台股量化交易工作站", Inches(6.5), Inches(6.88), Inches(6.5), Inches(0.5),
    size=11, color=MUTED)


# ── 儲存 ──────────────────────────────────────────────────────────────────────
out_path = "QUANT_AI_流程圖.pptx"
prs.save(out_path)
print(f"✅ 已生成：{out_path}")
print(f"   投影片數量：{len(prs.slides)} 張")
